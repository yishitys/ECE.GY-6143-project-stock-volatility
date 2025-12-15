#!/usr/bin/env python
"""
Generate a professional 10-minute PPT presentation for stock volatility prediction project
Focuses on: Formulation, Approach, Evaluation
Run: python generate_presentation.py
Output: Stock_Volatility_Prediction.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path
import re

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme (Professional Blue & Green)
COLOR_PRIMARY = RGBColor(46, 204, 113)    # Green
COLOR_ACCENT = RGBColor(52, 152, 219)    # Blue
COLOR_DARK = RGBColor(44, 62, 80)        # Dark
COLOR_LIGHT = RGBColor(236, 240, 241)    # Light gray
COLOR_RED = RGBColor(231, 76, 60)        # Red

def add_formatted_text(text_frame, text, font_size=Pt(18), font_color=COLOR_DARK, level=0):
    """Add text with bold formatting support. Use **text** for bold."""
    p = text_frame.add_paragraph() if text_frame.paragraphs[0].text else text_frame.paragraphs[0]
    p.level = level
    
    # Split by **bold** markers
    parts = re.split(r'\*\*([^*]+)\*\*', text)
    
    for i, part in enumerate(parts):
        if i == 0 and part:
            p.text = part
            p.font.size = font_size
            p.font.color.rgb = font_color
        elif part:
            if i % 2 == 1:  # Bold part
                run = p.add_run()
                run.text = part
                run.font.size = font_size
                run.font.bold = True
                run.font.color.rgb = font_color
            else:  # Normal part
                run = p.add_run()
                run.text = part
                run.font.size = font_size
                run.font.color.rgb = font_color
    
    return p

def add_title_slide(title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PRIMARY
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(255, 255, 255)

def add_content_slide(title, content_points, image_path=None):
    """Add a content slide with bullet points and optional image"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLOR_ACCENT
    title_shape.line.color.rgb = COLOR_ACCENT
    
    # Add title text (remove trailing colon if present)
    title_text = title.rstrip(':')
    title_frame = title_shape.text_frame
    title_frame.text = title_text
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Add content
    if image_path and Path(image_path).exists():
        # Left side: text, Right side: image
        text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.8), Inches(5.8))
        img_left = Inches(5.5)
        img_width = Inches(4)
    else:
        text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
        img_left = None
    
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(content_points):
        is_subheading = point.endswith(' ') or (not point.startswith("   ") and ":" in point and not point.startswith("•"))
        level = 0 if not point.startswith("  ") else 1
        
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.level = level
        display_text = point.lstrip()
        
        # Handle subheadings (bold them and remove trailing colon)
        if is_subheading and not display_text.startswith("•"):
            display_text = display_text.rstrip(":").rstrip() + ":"
            p.text = ""
            run = p.add_run()
            run.text = display_text
            run.font.size = Pt(18)
            run.font.bold = True
            run.font.color.rgb = COLOR_DARK
        else:
            # Use formatted text for bold keywords
            p.text = ""
            parts = re.split(r'\*\*([^*]+)\*\*', display_text)
            for j, part in enumerate(parts):
                if part:
                    run = p.add_run()
                    run.text = part
                    run.font.size = Pt(18)
                    run.font.color.rgb = COLOR_DARK
                    if j % 2 == 1:  # Bold part
                        run.font.bold = True
    
    # Add image if provided
    if image_path and Path(image_path).exists() and img_left:
        try:
            slide.shapes.add_picture(image_path, img_left, Inches(1.5), width=img_width)
        except:
            pass

def add_two_column_slide(title, left_content, right_content, image_path=None):
    """Add a two-column content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLOR_ACCENT
    title_shape.line.color.rgb = COLOR_ACCENT
    
    title_text = title.rstrip(':')
    title_frame = title_shape.text_frame
    title_frame.text = title_text
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.2), Inches(4.5), Inches(5.8))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    for i, point in enumerate(left_content):
        if i == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        
        display_text = point
        # Bold subheadings (lines with emoji + text)
        if point.strip() and not point.startswith("  ") and ":" in point:
            p.text = ""
            display_text = display_text.rstrip(":").rstrip() + ":"
            run = p.add_run()
            run.text = display_text
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = COLOR_DARK
        else:
            p.text = ""
            parts = re.split(r'\*\*([^*]+)\*\*', display_text)
            for j, part in enumerate(parts):
                if part:
                    run = p.add_run()
                    run.text = part
                    run.font.size = Pt(16)
                    run.font.color.rgb = COLOR_DARK
                    if j % 2 == 1:
                        run.font.bold = True
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.5), Inches(5.8))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    for i, point in enumerate(right_content):
        if i == 0:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()
        
        display_text = point
        # Bold subheadings
        if point.strip() and not point.startswith("  ") and ":" in point:
            p.text = ""
            display_text = display_text.rstrip(":").rstrip() + ":"
            run = p.add_run()
            run.text = display_text
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = COLOR_DARK
        else:
            p.text = ""
            parts = re.split(r'\*\*([^*]+)\*\*', display_text)
            for j, part in enumerate(parts):
                if part:
                    run = p.add_run()
                    run.text = part
                    run.font.size = Pt(16)
                    run.font.color.rgb = COLOR_DARK
                    if j % 2 == 1:
                        run.font.bold = True

# ============================================================================
# SLIDE 1: TITLE SLIDE
# ============================================================================
add_title_slide(
    "Stock Volatility Prediction",
    "Leveraging Social Media Discussion Data with Sentiment Embeddings"
)

# ============================================================================
# SLIDE 2: FORMULATION - Problem Statement
# ============================================================================
add_content_slide(
    "FORMULATION Problem Statement",
    [
        "❓ **Research Question**",
        "   Can **Reddit sentiment** help predict **stock volatility**",
        "   beyond **technical indicators** alone?",
        "",
        "💡 **Motivation**",
        "   • Market efficiency requires **all information**",
        "   • **Real-time, crowd-sourced** sentiment signals",
        "   • Reddit communities drive market behavior",
        "   • Challenge: **Extract semantic meaning**, not just volume"
    ]
)

# ============================================================================
# SLIDE 3: FORMULATION - Research Questions
# ============================================================================
add_content_slide(
    "FORMULATION 4 Research Questions",
    [
        "**RQ1**: Does Reddit volume correlate with volatility?",
        "   → Preliminary: **Weak (r=0.18)** → Need semantic features",
        "",
        "**RQ2**: Can text embeddings capture sentiment better?",
        "   → Hypothesis: **SBERT 6x better** than raw counts",
        "",
        "**RQ3**: Technical vs social data importance?",
        "   → Hypothesis: **Technical (70%)** + social (5-10%)",
        "",
        "**RQ4**: Can we achieve useful trading accuracy?",
        "   → Target: **>80% directional accuracy** next-hour"
    ]
)

# ============================================================================
# SLIDE 4: Data Overview
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
title_shape.fill.solid()
title_shape.fill.fore_color.rgb = COLOR_ACCENT
title_shape.line.color.rgb = COLOR_ACCENT
title_frame = title_shape.text_frame
title_frame.text = "Dataset: 2021 GME Trading Year"
p = title_frame.paragraphs[0]
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.8), Inches(5.8))
text_frame = text_box.text_frame
text_frame.word_wrap = True
points = [
    "📊 Scope:",
    "   • 2,001 hourly observations",
    "   • 168,158 Reddit posts",
    "   • Stock OHLCV data",
    "",
    "🎯 Target Variable:",
    "   • Hourly volatility",
    "   • Highly skewed (88% zeros)",
    "",
    "⚡ Temporal Alignment:",
    "   • UTC → EST conversion",
    "   • Hourly aggregation"
]
for i, point in enumerate(points):
    if i == 0:
        p = text_frame.paragraphs[0]
    else:
        p = text_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_DARK
    if point.startswith("•") or point.startswith("   •"):
        p.level = 1

if Path("data_visualization_diagrams/07_data_overview.png").exists():
    slide.shapes.add_picture("data_visualization_diagrams/07_data_overview.png", 
                             Inches(5.5), Inches(1.5), width=Inches(4))

# ============================================================================
# SLIDE 5: APPROACH - Technical Approach Overview
# ============================================================================
add_two_column_slide(
    "APPROACH End-to-End Pipeline",
    [
        "1️⃣ **Data Processing**",
        "   • **Align** Reddit/stock data",
        "   • **Clean** text (lowercase, remove URLs)",
        "",
        "2️⃣ **Feature Engineering**",
        "   • **384D embeddings** (SBERT)",
        "   • **50+ indicators**",
        "   • **4 Reddit stats**",
        "",
        "3️⃣ **Modeling**",
        "   • **XGBoost** (main)",
        "   • **LSTM** (baseline)",
        "   • **Baseline** models",
    ],
    [
        "4️⃣ **Evaluation**",
        "   • **Time-series CV**, no leakage",
        "   • **70/15/15** split",
        "   • **Multiple metrics**",
        "     - **R²** (variance)",
        "     - **MAE** (magnitude)",
        "     - **Dir. Acc** (signal)",
        "",
        "5️⃣ **Ablation Studies**",
        "   • Model **variants**",
        "   • **Feature importance**"
    ]
)

# ============================================================================
# SLIDE 6: APPROACH - Feature Engineering Strategy
# ============================================================================
add_content_slide(
    "APPROACH Feature Engineering Why 444 Features",
    [
        "🔴 **Technical Indicators** (50 features, 11.3%)",
        "   • Historical volatility, returns, **RSI, MACD, Bollinger Bands**",
        "   • Moving averages → **Baseline predictors**",
        "",
        "🟠 **Reddit Statistics** (4 features, 2.2%)",
        "   • Posts/hour, comments, unique authors",
        "   • **Too sparse** → Raw counts **weak (r=0.18)**",
        "",
        "🟡 **Text Embeddings** (384 features, 86.5%) **KEY INNOVATION**",
        "   • **All-MiniLM-L6-v2** (Sentence-BERT)",
        "   • Captures: **bullish/bearish/uncertainty**",
        "   • **6x better** than raw counts (0.30 vs 0.05)",
        "   • Challenge: **30+ hours compute**"
    ]
)

# ============================================================================
# SLIDE 7: APPROACH - Model Architecture
# ============================================================================
add_two_column_slide(
    "APPROACH Model Architectures",
    [
        "🟢 **XGBoost** (Chosen)",
        "   • **Depth**: 7",
        "   • **Learning rate**: 0.1",
        "   • **N_estimators**: 200",
        "   • Handles **heterogeneous features**",
        "   • **Feature importance** output",
        "   • **Interpretable**"
    ],
    [
        "🔵 **LSTM** (Baseline)",
        "   • **2 layers × 64 units**",
        "   • **24-hour sequences**",
        "   • **Dropout 0.2**",
        "   • Time-series capable",
        "",
        "⚪ **Baselines**",
        "   • **Persist**: Last = next",
        "   • **Mean**: Historical avg"
    ]
)

# ============================================================================
# SLIDE 8: EVALUATION - Evaluation Methodology
# ============================================================================
add_content_slide(
    "EVALUATION Rigorous Evaluation Strategy",
    [
        "✅ **Time-Series Cross-Validation** (Critical)",
        "   • Split: **70% train** (1,400h) → **15% val** → **15% test** (300h)",
        "   • **NO data leakage**, temporal ordering preserved",
        "",
        "📊 **Evaluation Metrics**",
        "   • **R²**: % variance explained (~27% ceiling)",
        "   • **MAE**: Absolute error (robust outliers)",
        "   • **Directional Accuracy**: Trading signal",
        "   • **Correlation**: Pearson r",
        "",
        "🔍 **Ablation Studies**",
        "   • Model 1: Full (technical + Reddit + embeddings)",
        "   • Model 2: **No text** ← **BEST**",
        "   • Model 3: Tech only (baseline)"
    ]
)

# ============================================================================
# SLIDE 9: RESULTS - Model Comparison (R² and MAE)
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
title_shape.fill.solid()
title_shape.fill.fore_color.rgb = COLOR_ACCENT
title_shape.line.color.rgb = COLOR_ACCENT
title_frame = title_shape.text_frame
title_frame.text = "RESULTS: Model Performance Comparison"
p = title_frame.paragraphs[0]
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

# Two-column layout with images on each side
if Path("data_visualization_diagrams/01_model_r2_comparison.png").exists():
    slide.shapes.add_picture("data_visualization_diagrams/01_model_r2_comparison.png",
                             Inches(0.3), Inches(1.0), width=Inches(4.6))
if Path("data_visualization_diagrams/02_model_mae_comparison.png").exists():
    slide.shapes.add_picture("data_visualization_diagrams/02_model_mae_comparison.png",
                             Inches(5.1), Inches(1.0), width=Inches(4.6))

# Add text box for key insights
text_box = slide.shapes.add_textbox(Inches(0.3), Inches(5.2), Inches(9.4), Inches(2))
text_frame = text_box.text_frame
text_frame.word_wrap = True

p = text_frame.paragraphs[0]
p.text = "🥇 Winner: XGBoost-no_text (R² = 0.2706, MAE = 0.0081, Dir. Acc. = 85.14%) | Key Insight: Text embeddings ADD NOISE, not signal"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = COLOR_RED

# ============================================================================
# SLIDE 10: RESULTS - Feature Importance & Category Contribution
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
title_shape.fill.solid()
title_shape.fill.fore_color.rgb = COLOR_ACCENT
title_shape.line.color.rgb = COLOR_ACCENT
title_frame = title_shape.text_frame
title_frame.text = "RESULTS: What Drives Volatility? (Feature Importance)"
p = title_frame.paragraphs[0]
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

# Left image: Feature importance
if Path("data_visualization_diagrams/04_feature_importance.png").exists():
    slide.shapes.add_picture("data_visualization_diagrams/04_feature_importance.png",
                             Inches(0.3), Inches(1.0), width=Inches(4.8))

# Right image: Feature category contribution
if Path("data_visualization_diagrams/05_feature_category_contribution.png").exists():
    slide.shapes.add_picture("data_visualization_diagrams/05_feature_category_contribution.png",
                             Inches(5.1), Inches(1.0), width=Inches(4.6))

# Add insights at bottom
text_box = slide.shapes.add_textbox(Inches(0.3), Inches(5.2), Inches(9.4), Inches(2))
text_frame = text_box.text_frame
text_frame.word_wrap = True

p = text_frame.paragraphs[0]
p.text = "🔴 Critical Finding: "
p.font.size = Pt(14)
p.font.color.rgb = COLOR_DARK

run = p.add_run()
run.text = "TECHNICAL INDICATORS DOMINATE (70%)"
run.font.bold = True
run.font.color.rgb = COLOR_RED

p = text_frame.add_paragraph()
p.text = "Top Features: Historical Volatility (12.5%) → Previous Returns (11.2%) → RSI (9.8%) | Reddit Discussion: Only 5-10% Additional Value"
p.font.size = Pt(13)
p.font.color.rgb = COLOR_DARK

# ============================================================================
# SLIDE 11: RESULTS - Directional Accuracy & Regime Analysis
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
title_shape.fill.solid()
title_shape.fill.fore_color.rgb = COLOR_ACCENT
title_shape.line.color.rgb = COLOR_ACCENT
title_frame = title_shape.text_frame
title_frame.text = "RESULTS: When Does Our Model Work Best?"
p = title_frame.paragraphs[0]
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

# Left image: Directional accuracy
if Path("data_visualization_diagrams/03_model_directional_accuracy.png").exists():
    slide.shapes.add_picture("data_visualization_diagrams/03_model_directional_accuracy.png",
                             Inches(0.3), Inches(1.0), width=Inches(4.8))

# Right image: By regime
if Path("data_visualization_diagrams/09_directional_by_regime.png").exists():
    slide.shapes.add_picture("data_visualization_diagrams/09_directional_by_regime.png",
                             Inches(5.1), Inches(1.0), width=Inches(4.6))

# Add insights
text_box = slide.shapes.add_textbox(Inches(0.3), Inches(5.2), Inches(9.4), Inches(2))
text_frame = text_box.text_frame
text_frame.word_wrap = True

p = text_frame.paragraphs[0]
p.text = "📈 Conditional Performance: Low Vol (92% ✅) → Medium (78% 🟡) → High (64% ❌) | "
p.font.size = Pt(13)
p.font.color.rgb = COLOR_DARK

run = p.add_run()
run.text = "Model excels in normal conditions, fails on exogenous shocks"
run.font.bold = True
run.font.color.rgb = COLOR_RED

# ============================================================================
# SLIDE 12: Feature Comparison & Time Series Analysis
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
title_shape.fill.solid()
title_shape.fill.fore_color.rgb = COLOR_ACCENT
title_shape.line.color.rgb = COLOR_ACCENT
title_frame = title_shape.text_frame
title_frame.text = "RESULTS: Reddit vs Technical + Time Series Predictions"
p = title_frame.paragraphs[0]
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

# Left: Reddit vs Technical comparison
if Path("data_visualization_diagrams/06_reddit_vs_technical.png").exists():
    slide.shapes.add_picture("data_visualization_diagrams/06_reddit_vs_technical.png",
                             Inches(0.3), Inches(1.0), width=Inches(4.8))

# Right: Time series predictions
if Path("data_visualization_diagrams/08_timeseries_predictions.png").exists():
    slide.shapes.add_picture("data_visualization_diagrams/08_timeseries_predictions.png",
                             Inches(5.1), Inches(1.0), width=Inches(4.6))

# Add insights
text_box = slide.shapes.add_textbox(Inches(0.3), Inches(5.2), Inches(9.4), Inches(2))
text_frame = text_box.text_frame
text_frame.word_wrap = True

p = text_frame.paragraphs[0]
p.text = "Left: Reddit embeddings (0.30) "
p.font.size = Pt(13)
p.font.color.rgb = COLOR_DARK

run = p.add_run()
run.text = "6x better"
run.font.bold = True

p = text_frame.add_paragraph()
p.text = "than raw counts (0.05), but technical still dominates | Right: Predicted vs actual volatility show strong correlation especially in calm periods"
p.font.size = Pt(13)
p.font.color.rgb = COLOR_DARK

# ============================================================================
# SLIDE 13: Key Findings & Implications
# ============================================================================
add_content_slide(
    "KEY FINDINGS Implications",
    [
        "🎯 **Finding 1**: **Embeddings > Raw Counts**",
        "   • Semantic meaning **6x more predictive**",
        "   • **Bullish/bearish** language matters",
        "",
        "🎯 **Finding 2**: **Technical Dominates**",
        "   • **70% feature importance** from technical",
        "   • Information **already priced in** (EMH)",
        "",
        "🎯 **Finding 3**: **Strong Directional Signal**",
        "   • **85% accuracy** (vs 50% random)",
        "   • **Actionable** for hedging",
        "",
        "💼 **Practical Implication**",
        "   • Use as **supplementary signal**",
        "   • Combine with **risk management**"
    ]
)

# ============================================================================
# SLIDE 14: Summary Metrics & Evaluation Summary
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
title_shape.fill.solid()
title_shape.fill.fore_color.rgb = COLOR_ACCENT
title_shape.line.color.rgb = COLOR_ACCENT
title_frame = title_shape.text_frame
title_frame.text = "RESULTS: Performance Summary & Evaluation"
p = title_frame.paragraphs[0]
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

# Images side by side
if Path("data_visualization_diagrams/10_metrics_summary_table.png").exists():
    slide.shapes.add_picture("data_visualization_diagrams/10_metrics_summary_table.png",
                             Inches(0.3), Inches(1.0), width=Inches(9.4), height=Inches(5.5))

# ============================================================================
# SLIDE 15: Limitations & Future Work
# ============================================================================
add_two_column_slide(
    "LIMITATIONS FUTURE WORK",
    [
        "🔴 **Limitations**",
        "",
        "• **Single stock** (GME)",
        "  → **High volatility**",
        "",
        "• **Single year** (2021)",
        "  → **Bull market**",
        "",
        "• **No transaction costs**",
        "  → **Slippage** ignored",
        "",
        "• **Cold start** problem",
        "  → New discussions"
    ],
    [
        "🟢 **Future Work**",
        "",
        "• **Multi-stock** portfolio",
        "  → Test **generalization**",
        "",
        "• **Temporal embeddings**",
        "  → **Sentiment drift**",
        "",
        "• **Causal inference**",
        "  → Reddit **causes** moves?",
        "",
        "• **Real trading** backtest",
        "  → **Profitability** test"
    ]
)

# ============================================================================
# SLIDE 16: CONCLUSION
# ============================================================================
add_title_slide(
    "Conclusion",
    "Social Media Signals Can Enhance (But Not Replace) Market Analysis"
)

# Save presentation
output_file = Path("Stock_Volatility_Prediction.pptx")
prs.save(output_file)
print(f"✅ PPT generated successfully!")
print(f"📁 Saved as: {output_file.absolute()}")
print(f"\n📊 Presentation includes:")
print("   • 16 slides with all 10 visualizations")
print("   • Focus on Formulation (Problem + RQ)")
print("   • Focus on Approach (Technical strategy)")
print("   • Focus on Evaluation (Results + all charts)")
print("   • Professional color scheme (Blue/Green)")
print("   • Bold highlights on key insights")

